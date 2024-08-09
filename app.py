import os
import uuid
import logging
import schedule
import time
from threading import Thread
from flask import Flask, request, jsonify, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY')

# Directory to save the generated presentations
PRESENTATION_DIR = 'presentations'
if not os.path.exists(PRESENTATION_DIR):
    os.makedirs(PRESENTATION_DIR)

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def delete_old_presentations():
    logger.info("Running scheduled cleanup task")
    for filename in os.listdir(PRESENTATION_DIR):
        if filename.endswith('.pptx'):
            file_path = os.path.join(PRESENTATION_DIR, filename)
            logger.info(f"Attempting to delete file: {file_path}")
            try:
                os.remove(file_path)
                logger.info(f"Deleted {file_path}")
            except Exception as e:
                logger.error(f"Error deleting file {file_path}: {e}")

# Schedule the cleanup task to run every 2 minutes
schedule.every(2).minutes.do(delete_old_presentations)

def run_scheduler():
    logger.info("Scheduler started")
    while True:
        schedule.run_pending()
        time.sleep(1)

@app.route('/generate_presentation', methods=['POST'])
def generate_presentation():
    try:
        data = request.get_json()

        if not data or 'slides_content' not in data:
            return jsonify({"error": "Invalid input"}), 400

        slides_content = data['slides_content']

        # Create a presentation object
        # Load the template presentation
        prs = Presentation()



        # Set aspect ratio to 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Loop through the content and create slides
        for slide_content in slides_content:
            slide_layout = prs.slide_layouts[5]  # Use a blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            title = slide.shapes.title
            if not title:
                title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
            title.text = slide_content.get("title", "No Title")

            # Format title
            title_text_frame = title.text_frame
            for paragraph in title_text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.size = Pt(32)  # Adjust the font size as needed
                paragraph.font.bold = True

            # Add content
            content_text = slide_content.get("content", "No Content")
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), prs.slide_height - Inches(2))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            # Split the content text by line breaks and add each line as a paragraph
            for line in content_text.split('\n'):
                p = text_frame.add_paragraph()
                p.text = line
                p.level = 0  # Ensure no bullet points
                p.font.size = Pt(18)  # Adjust the font size as needed
                p.font.bold = False
                p.space_before = Pt(10)  # Adjust the space before each paragraph
                p.space_after = Pt(10)  # Adjust the space after each paragraph

        # Generate a unique filename
        unique_filename = f"{uuid.uuid4()}.pptx"
        file_path = os.path.join(PRESENTATION_DIR, unique_filename)

        # Save the presentation to a file
        prs.save(file_path)

        logger.info(f"Presentation created: {file_path}")

        return jsonify({"message": "Presentation created successfully", "file_url": f"/download/{unique_filename}"}), 200
    except Exception as e:
        logger.error(f"Error creating presentation: {e}")
        return jsonify({"error": "Internal Server Error"}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    try:
        return send_from_directory(PRESENTATION_DIR, filename)
    except Exception as e:
        logger.error(f"Error downloading file: {e}")
        return jsonify({"error": "File not found"}), 404

if __name__ == '__main__':
    scheduler_thread = Thread(target=run_scheduler)
    scheduler_thread.daemon = True
    scheduler_thread.start()
    app.run(debug=True)
